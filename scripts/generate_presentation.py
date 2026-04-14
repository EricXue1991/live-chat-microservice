#!/usr/bin/env python3
"""Generate a polished LiveChat PowerPoint presentation.

Install deps:
    pip install -r scripts/requirements-presentation.txt
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
C_NAVY   = RGBColor(0x1B, 0x3A, 0x6B)   # deep blue
C_TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # accent teal
C_TEAL2  = RGBColor(0x48, 0xCA, 0xE4)   # lighter teal
C_ORANGE = RGBColor(0xFF, 0x6B, 0x35)   # accent orange (highlights)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xE8, 0xF4, 0xF8)   # very light blue-grey (table rows)
C_DGRAY  = RGBColor(0x4A, 0x5A, 0x6A)   # body text on light bg
C_GREEN  = RGBColor(0x2D, 0xC6, 0x53)   # positive results
C_RED    = RGBColor(0xE6, 0x3F, 0x3F)   # negative / error

W = Inches(13.333)
H = Inches(7.5)
FIGS = Path(__file__).resolve().parent.parent / "report" / "figures"


# ── Low-level helpers ────────────────────────────────────────────────────────

def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def solid_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_font(run_or_para, size_pt, bold=False, color: RGBColor = None, italic=False):
    f = run_or_para.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    if color:
        f.color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    solid_fill(shape, fill_color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_textbox(slide, left, top, width, height, text, size_pt=18,
                bold=False, color: RGBColor = C_WHITE, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold=bold, color=color, italic=italic)
    return txb


def add_image_safe(slide, img_path, left, top, width, height):
    p = Path(img_path)
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width, height)
        return True
    return False


# ── Background helper ────────────────────────────────────────────────────────

def dark_bg(slide):
    """Full dark navy background."""
    add_rect(slide, 0, 0, W, H, C_DARK)


def light_bg(slide):
    """White/very-light background for content slides."""
    add_rect(slide, 0, 0, W, H, C_WHITE)


def header_bar(slide, text, subtitle=None):
    """Top accent bar + slide title."""
    add_rect(slide, 0, 0, W, Inches(1.35), C_NAVY)
    add_rect(slide, 0, Inches(1.35), Inches(0.06), H - Inches(1.35), C_TEAL)
    add_textbox(slide, Inches(0.3), Inches(0.12), W - Inches(0.5), Inches(0.85),
                text, size_pt=32, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.3), Inches(0.92), W - Inches(0.5), Inches(0.42),
                    subtitle, size_pt=15, color=C_TEAL2, italic=True)


def teal_dot(slide, cx, cy, r=Inches(0.07)):
    """Small decorative dot."""
    shape = slide.shapes.add_shape(9, cx - r, cy - r, r * 2, r * 2)  # oval
    solid_fill(shape, C_TEAL)
    shape.line.fill.background()


# ── Bullet helper ─────────────────────────────────────────────────────────────

def add_bullets(slide, left, top, width, height,
                items,           # list of (text, level, size, bold, color)
                spacing_pt=6):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, level, size, bold, color = item, 0, 18, False, C_DGRAY
        else:
            text, level, size, bold, color = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(spacing_pt) if level == 0 else Pt(2)
        indent = Inches(0.2 * level)
        bullet_char = "▸ " if level == 0 else "  • "
        run = p.add_run()
        run.text = bullet_char + text
        set_font(run, size, bold=bold, color=color)


# ── Table helper ─────────────────────────────────────────────────────────────

def add_table(slide, left, top, width, rows, cols, headers, data,
              col_widths=None, font_size=13):
    tbl = slide.shapes.add_table(len(data) + 1, cols, left, top, width, Inches(0.4)).table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(width * w / total)

    def cell_fill(cell, color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr.set('val', rgb_hex(color))

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell_fill(cell, C_NAVY)
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = h
        set_font(run, font_size, bold=True, color=C_WHITE)

    # Data rows
    for i, row in enumerate(data):
        bg = C_LGRAY if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell_fill(cell, bg)
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            if isinstance(val, tuple):
                run.text, extra_color = val[0], val[1]
                set_font(run, font_size, bold=True, color=extra_color)
            else:
                run.text = str(val)
                set_font(run, font_size, color=C_DGRAY)

    return tbl


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Full dark background
    add_rect(slide, 0, 0, W, H, C_DARK)
    # Large teal accent strip on left
    add_rect(slide, 0, 0, Inches(0.45), H, C_TEAL)
    # Bottom accent line
    add_rect(slide, Inches(0.45), H - Inches(0.12), W - Inches(0.45), Inches(0.12), C_ORANGE)

    # Main title
    add_textbox(slide, Inches(1.0), Inches(1.4), Inches(11), Inches(1.6),
                "LiveChat Microservice", size_pt=54, bold=True, color=C_WHITE)
    # Teal underline
    add_rect(slide, Inches(1.0), Inches(3.05), Inches(5.5), Inches(0.06), C_TEAL)

    # Subtitle
    add_textbox(slide, Inches(1.0), Inches(3.25), Inches(10), Inches(0.7),
                "Distributed Real-Time Chat · Reaction Aggregation · AWS ECS Experiments",
                size_pt=20, color=C_TEAL2)

    # Details
    add_textbox(slide, Inches(1.0), Inches(4.2), Inches(9), Inches(0.5),
                "CS 6650  ·  Locust Load Testing  ·  PostgreSQL · DynamoDB · Redis · Kafka · SNS/SQS",
                size_pt=15, color=RGBColor(0xAA, 0xCC, 0xDD))

    # Decorative dots
    for y in [Inches(5.5), Inches(5.9), Inches(6.3)]:
        teal_dot(slide, Inches(1.3), y)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Agenda")

    items = [
        ("System Overview & Architecture", 0, 20, True,  C_NAVY),
        ("Tech Stack",                      0, 20, True,  C_NAVY),
        ("Experiment 1 — Horizontal Scale-Out",          0, 20, True,  C_NAVY),
        ("Experiment 2 — Hot Room vs Multi-Room",        0, 20, True,  C_NAVY),
        ("Experiment 3 — Sync vs Async Reactions",       0, 20, True,  C_NAVY),
        ("Experiment 4 — WebSocket vs HTTP Polling",     0, 20, True,  C_NAVY),
        ("Key Findings & Takeaways",         0, 20, True,  C_NAVY),
    ]

    step_x = Inches(1.2)
    step_y = Inches(1.6)
    for idx, (text, level, size, bold, color) in enumerate(items):
        y = step_y + idx * Inches(0.72)
        # Number badge
        badge = slide.shapes.add_shape(9, step_x - Inches(0.42), y - Inches(0.05),
                                       Inches(0.38), Inches(0.38))
        solid_fill(badge, C_TEAL if idx < 2 else C_NAVY)
        badge.line.fill.background()
        add_textbox(slide, step_x - Inches(0.44), y - Inches(0.07), Inches(0.42), Inches(0.42),
                    str(idx + 1), size_pt=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, step_x + Inches(0.05), y, Inches(10), Inches(0.45),
                    text, size_pt=size, bold=bold, color=color)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "System Architecture", "Go backend · React frontend · AWS-native data layer")

    # Draw architecture boxes
    box_data = [
        ("React\nFrontend\n:3000",      Inches(0.5),  Inches(2.2), Inches(1.6), Inches(1.4), C_TEAL),
        ("ALB\n(port 80/443)",          Inches(2.4),  Inches(2.2), Inches(1.6), Inches(1.4), C_NAVY),
        ("Go API\nReplicas\n(ECS)",     Inches(4.3),  Inches(2.2), Inches(1.6), Inches(1.4), C_ORANGE),
        ("PostgreSQL\nUsers / Rooms",   Inches(6.5),  Inches(1.2), Inches(1.8), Inches(1.1), C_NAVY),
        ("DynamoDB\nMessages &\nReactions", Inches(6.5), Inches(2.5), Inches(1.8), Inches(1.2), C_NAVY),
        ("Redis\nCache / Rate\nLimit / Presence", Inches(6.5), Inches(3.9), Inches(1.8), Inches(1.2), C_NAVY),
        ("S3\nAttachments",             Inches(8.6),  Inches(1.2), Inches(1.7), Inches(1.0), C_TEAL),
        ("Kafka\nAnalytics\nStream",    Inches(8.6),  Inches(2.4), Inches(1.7), Inches(1.1), C_TEAL),
        ("SNS → SQS\nCross-replica\nWS Fan-out", Inches(8.6), Inches(3.65), Inches(1.7), Inches(1.3), C_TEAL),
    ]

    for label, l, t, w2, h2, color in box_data:
        r = slide.shapes.add_shape(1, l, t, w2, h2)
        solid_fill(r, color)
        r.line.color.rgb = C_WHITE
        tf = r.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        set_font(run, 10, bold=True, color=C_WHITE)
        # vertical center hack via space_before
        tf.paragraphs[0].space_before = Pt(4)

    # Arrows (simple lines via connectors)
    def arrow(slide, x1, y1, x2, y2):
        from pptx.util import Emu
        connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
        connector.line.color.rgb = C_TEAL
        connector.line.width = Pt(1.5)

    arrow(slide, Inches(2.1), Inches(2.9), Inches(2.4), Inches(2.9))
    arrow(slide, Inches(4.0), Inches(2.9), Inches(4.3), Inches(2.9))
    arrow(slide, Inches(5.9), Inches(2.6), Inches(6.5), Inches(1.7))
    arrow(slide, Inches(5.9), Inches(2.9), Inches(6.5), Inches(3.0))
    arrow(slide, Inches(5.9), Inches(3.2), Inches(6.5), Inches(4.3))
    arrow(slide, Inches(8.3), Inches(1.6), Inches(8.6), Inches(1.6))
    arrow(slide, Inches(8.3), Inches(2.9), Inches(8.6), Inches(2.9))
    arrow(slide, Inches(8.3), Inches(4.2), Inches(8.6), Inches(4.2))

    # DynamoDB key design note
    add_rect(slide, Inches(0.3), Inches(5.6), Inches(12.5), Inches(1.6), C_LGRAY)
    add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12.0), Inches(1.4),
                "DynamoDB key design: PK = room_id  ·  SK = timestamp#messageId  →  "
                "efficient time-range queries + write distribution across partitions",
                size_pt=14, color=C_NAVY, bold=False)


def slide_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Tech Stack")

    cols_data = [
        ("Frontend",  ["React + Vite", "TypeScript", "Tailwind CSS", "useWebSocket hook"]),
        ("Backend",   ["Go (single binary)", "JWT auth (stateless)", "WebSocket (gorilla/ws)", "Kafka producer"]),
        ("Storage",   ["PostgreSQL 16 (relational)", "DynamoDB on-demand (chat)", "Redis 7 (cache / rate-limit / presence)", "S3 (attachments)"]),
        ("Infra",     ["AWS ECS Fargate", "Application Load Balancer", "SNS → SQS (fan-out)", "Terraform (IaC)"]),
        ("Testing",   ["Locust (Python)", "ChatUser, ReactionHeavyUser", "PollingUser, WebSocketUser", "AWS & LocalStack runs"]),
    ]

    col_w = Inches(2.4)
    for ci, (title, items) in enumerate(cols_data):
        lx = Inches(0.3) + ci * col_w
        # header badge
        add_rect(slide, lx, Inches(1.5), col_w - Inches(0.1), Inches(0.5), C_NAVY)
        add_textbox(slide, lx, Inches(1.53), col_w - Inches(0.1), Inches(0.45),
                    title, size_pt=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        for ri, item in enumerate(items):
            ry = Inches(2.15) + ri * Inches(0.9)
            add_rect(slide, lx, ry, col_w - Inches(0.1), Inches(0.8),
                     C_LGRAY if ri % 2 == 0 else C_WHITE)
            add_textbox(slide, lx + Inches(0.1), ry + Inches(0.1),
                        col_w - Inches(0.2), Inches(0.7),
                        item, size_pt=13, color=C_DGRAY)


def slide_exp1_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 1 — Horizontal Scale-Out",
               "150 users · 120 s runs · ECS Fargate behind ALB · AWS us-west-2")

    headers = ["Replicas", "Throughput (req/s)", "Avg Latency (ms)", "p95 (ms)", "p99 (ms)", "Error Rate", "Scale Eff."]
    data = [
        ["1",  "38.2", "1289", "14 000", "22 000", ("4.22%", C_RED),   "100%"],
        ["2",  "55.2",  "347",    "120",  "9 500", ("3.8%",  C_RED),   ("72.3%", C_GREEN)],
        ["4",  "56.5",  "323",    "110",  "8 800", ("3.5%",  C_ORANGE), ("37.0%", C_DGRAY)],
        ["8",  ("57.0", C_GREEN), ("312", C_GREEN), "110",  "8 500", ("3.2%",  C_GREEN), ("18.7%", C_DGRAY)],
    ]
    add_table(slide, Inches(0.4), Inches(1.65), Inches(12.4), 5, 7, headers, data,
              col_widths=[1, 2, 2.2, 1.5, 1.5, 1.8, 1.5], font_size=14)

    # Key finding box
    add_rect(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.9), RGBColor(0xE8, 0xF8, 0xFF))
    add_textbox(slide, Inches(0.55), Inches(4.05), Inches(12.0), Inches(0.8),
                "▸ 1→2 replicas: +44% throughput, −73% avg latency  "
                "▸ Beyond 2 replicas: throughput plateaus ~57 req/s  "
                "▸ Bottleneck: shared backends (PostgreSQL login burst, DynamoDB, ALB)",
                size_pt=14, color=C_NAVY)


def slide_exp1_charts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 1 — Charts")

    imgs = [
        (FIGS / "exp1" / "exp1_throughput_vs_replicas.png",  Inches(0.3),  Inches(1.5), Inches(4.1), Inches(3.0)),
        (FIGS / "exp1" / "exp1_latency_vs_replicas.png",     Inches(4.6),  Inches(1.5), Inches(4.1), Inches(3.0)),
        (FIGS / "exp1" / "exp1_scaling_efficiency.png",      Inches(9.0),  Inches(1.5), Inches(4.0), Inches(3.0)),
        (FIGS / "exp1" / "exp1_throughput_timeseries.png",   Inches(0.3),  Inches(4.65), Inches(6.0), Inches(2.6)),
        (FIGS / "exp1" / "exp1_error_rate_vs_replicas.png",  Inches(6.7),  Inches(4.65), Inches(6.0), Inches(2.6)),
    ]
    captions = ["Throughput vs Replicas", "Latency vs Replicas",
                "Scaling Efficiency", "Throughput Time-series", "Error Rate"]
    for idx, (path, l, t, w2, h2) in enumerate(imgs):
        add_image_safe(slide, path, l, t, w2, h2)
        add_textbox(slide, l, t + h2, w2, Inches(0.3),
                    captions[idx], size_pt=10, color=C_DGRAY, align=PP_ALIGN.CENTER)


def slide_exp2_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 2 — Hot Room vs Multi-Room",
               "500 users · 180 s · DynamoDB partition design · LocalStack vs AWS")

    # Two side-by-side result tables
    # LocalStack
    add_rect(slide, Inches(0.3), Inches(1.6), Inches(6.0), Inches(0.45), C_NAVY)
    add_textbox(slide, Inches(0.3), Inches(1.62), Inches(6.0), Inches(0.4),
                "LocalStack Results (Application-layer effects dominate)",
                size_pt=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    headers2 = ["Metric", "Hot Room", "Multi Room"]
    data_local = [
        ["Throughput (req/s)", "298.2", ("278.5", C_RED)],
        ["Avg Latency (ms)",   ("994",  C_GREEN), "1 112"],
        ["p99 Latency – React (ms)", ("1 550", C_GREEN), "2 100"],
        ["Error Rate", "0%", "0%"],
    ]
    add_table(slide, Inches(0.3), Inches(2.1), Inches(6.0), 5, 3, headers2, data_local,
              col_widths=[2.5, 1.5, 1.5], font_size=13)

    # AWS
    add_rect(slide, Inches(6.9), Inches(1.6), Inches(6.1), Inches(0.45), C_TEAL)
    add_textbox(slide, Inches(6.9), Inches(1.62), Inches(6.1), Inches(0.4),
                "AWS Results (Real DynamoDB partition contention)",
                size_pt=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    data_aws = [
        ["Throughput (req/s)", "52.3", ("54.8", C_GREEN)],
        ["Avg Latency (ms)",   "186",  ("142",  C_GREEN)],
        ["p99 Latency (ms)",   "9 200", ("6 800", C_GREEN)],
        ["Error Rate",         ("4.1%", C_RED), ("3.6%", C_ORANGE)],
    ]
    add_table(slide, Inches(6.9), Inches(2.1), Inches(6.1), 5, 3, headers2, data_aws,
              col_widths=[2.5, 1.5, 1.5], font_size=13)

    # Insight
    add_rect(slide, Inches(0.3), Inches(4.75), Inches(12.7), Inches(1.05), RGBColor(0xFF, 0xF3, 0xE0))
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.85),
                "LocalStack: hot room wins — single mutex entry, less lock contention at app layer.\n"
                "AWS: multi-room wins — real DynamoDB parallelises writes across physical partitions, "
                "reducing hot-partition throttling (4.1% vs 3.6% errors).",
                size_pt=13, color=C_NAVY)


def slide_exp2_charts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 2 — Charts")

    imgs = [
        (FIGS / "exp2"     / "exp2_throughput.png",          Inches(0.3),  Inches(1.5), Inches(4.0), Inches(2.8), "Local – Throughput"),
        (FIGS / "exp2"     / "exp2_latency.png",             Inches(4.5),  Inches(1.5), Inches(4.0), Inches(2.8), "Local – Latency"),
        (FIGS / "exp2"     / "exp2_throughput_timeseries.png",Inches(8.9),  Inches(1.5), Inches(4.1), Inches(2.8), "Local – Time-series"),
        (FIGS / "exp2_aws" / "exp2_throughput.png",           Inches(0.3),  Inches(4.5), Inches(4.0), Inches(2.8), "AWS – Throughput"),
        (FIGS / "exp2_aws" / "exp2_latency.png",              Inches(4.5),  Inches(4.5), Inches(4.0), Inches(2.8), "AWS – Latency"),
        (FIGS / "exp2_aws" / "exp2_throughput_timeseries.png",Inches(8.9),  Inches(4.5), Inches(4.1), Inches(2.8), "AWS – Time-series"),
    ]
    for path, l, t, w2, h2, cap in imgs:
        add_image_safe(slide, path, l, t, w2, h2)
        add_textbox(slide, l, t + h2 + Inches(0.02), w2, Inches(0.3),
                    cap, size_pt=10, color=C_DGRAY, align=PP_ALIGN.CENTER)


def slide_exp3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 3 — Sync vs Async Reactions",
               "500 users · 180 s · REACTION_MODE=sync/async · SQS batch aggregation")

    # Left: explanation
    add_textbox(slide, Inches(0.4), Inches(1.55), Inches(5.5), Inches(0.35),
                "Sync Mode", size_pt=16, bold=True, color=C_NAVY)
    add_bullets(slide, Inches(0.4), Inches(1.95), Inches(5.5), Inches(1.5), [
        ("Each reaction → immediate DynamoDB counter write", 0, 13, False, C_DGRAY),
        ("High WCU consumption on hot counters", 0, 13, False, C_DGRAY),
        ("Simple, predictable latency", 0, 13, False, C_DGRAY),
    ])

    add_rect(slide, Inches(6.0), Inches(1.55), Inches(0.05), Inches(2.0), C_TEAL)

    add_textbox(slide, Inches(6.3), Inches(1.55), Inches(5.5), Inches(0.35),
                "Async Mode (SQS Batch)", size_pt=16, bold=True, color=C_TEAL)
    add_bullets(slide, Inches(6.3), Inches(1.95), Inches(5.5), Inches(1.5), [
        ("Reactions → SQS queue; batch consumer aggregates ~10 s", 0, 13, False, C_DGRAY),
        ("Reduces WCU spikes; back-pressure via queue depth", 0, 13, False, C_DGRAY),
        ("Queue visible at /api/status (reaction_queue_visible)", 0, 13, False, C_DGRAY),
    ])

    # Charts
    imgs = [
        (FIGS / "exp3"     / "exp3_reaction_throughput.png",  Inches(0.3),  Inches(3.9), Inches(4.0), Inches(3.2), "Local – Throughput"),
        (FIGS / "exp3"     / "exp3_reaction_latency.png",     Inches(4.5),  Inches(3.9), Inches(4.0), Inches(3.2), "Local – Latency"),
        (FIGS / "exp3_aws" / "exp3_summary_throughput_latency.png", Inches(8.9), Inches(3.9), Inches(4.1), Inches(3.2), "AWS – Summary"),
    ]
    for path, l, t, w2, h2, cap in imgs:
        add_image_safe(slide, path, l, t, w2, h2)
        add_textbox(slide, l, t + h2 + Inches(0.02), w2, Inches(0.3),
                    cap, size_pt=10, color=C_DGRAY, align=PP_ALIGN.CENTER)


def slide_exp4_bug(prs):
    """The bug discovery story for Exp 4."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 4 — WebSocket vs HTTP Polling: The Bug",
               "50 users · WebSocket latency was WORSE than polling before the fix")

    # Bug results
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(5.8), Inches(0.45), C_RED)
    add_textbox(slide, Inches(0.3), Inches(1.57), Inches(5.8), Inches(0.4),
                "Pre-fix Results (BUGGY)", size_pt=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    headers_e = ["Metric", "HTTP Polling", "WebSocket"]
    data_bug = [
        ["p50 (ms)",  "1 000", ("17 000", C_RED)],
        ["p95 (ms)",  "3 200", ("33 000", C_RED)],
        ["p99 (ms)",  "4 600", ("37 000", C_RED)],
        ["Average",   "1 261", ("17 160", C_RED)],
    ]
    add_table(slide, Inches(0.3), Inches(2.05), Inches(5.8), 5, 3, headers_e, data_bug,
              col_widths=[2, 1.8, 1.8], font_size=13)

    # Root cause
    add_rect(slide, Inches(0.3), Inches(4.45), Inches(5.8), Inches(0.45), RGBColor(0xFF, 0xE5, 0xE5))
    add_textbox(slide, Inches(0.5), Inches(4.48), Inches(5.6), Inches(0.4),
                "Root Cause", size_pt=13, bold=True, color=C_RED)
    add_textbox(slide, Inches(0.4), Inches(4.95), Inches(5.6), Inches(1.5),
                "ALL WebSocket delivery routed through SNS→SQS\n"
                "(WaitTimeSeconds=20), even for clients on the SAME replica.\n"
                "Missing: direct hub.BroadcastToRoom() call.",
                size_pt=13, color=C_DGRAY)

    # Fix description
    add_rect(slide, Inches(6.4), Inches(1.55), Inches(6.5), Inches(0.45), C_GREEN)
    add_textbox(slide, Inches(6.4), Inches(1.57), Inches(6.5), Inches(0.4),
                "The Fix", size_pt=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(6.5), Inches(2.05), Inches(6.2), Inches(2.4), [
        ("models/models.go — Added SourceID to BroadcastMessage", 0, 13, False, C_DGRAY),
        ("ws/hub.go — Added replica ID, skip own SQS messages", 0, 13, False, C_DGRAY),
        ("chat/handler.go — Direct hub.BroadcastToRoom() call", 0, 13, False, C_DGRAY),
        ("cmd/server/main.go — Wire hub into chat handler", 0, 13, False, C_DGRAY),
    ])
    # Flow diagram (text)
    add_rect(slide, Inches(6.4), Inches(3.95), Inches(6.5), Inches(2.8), C_LGRAY)
    add_textbox(slide, Inches(6.6), Inches(4.05), Inches(6.1), Inches(2.6),
                "Fixed delivery path:\n\n"
                "POST /api/messages\n"
                "  ├─▸ hub.BroadcastToRoom()  ←  direct, ~20 ms\n"
                "  └─▸ SNS → SQS  ←  cross-replica fan-out (OK)\n\n"
                "SourceID prevents duplicate delivery",
                size_pt=13, color=C_NAVY)


def slide_exp4_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    light_bg(slide)
    header_bar(slide, "Experiment 4 — Post-fix Results",
               "Local vs AWS · WebSocket dramatically outperforms HTTP polling")

    # Local
    add_rect(slide, Inches(0.3), Inches(1.55), Inches(5.8), Inches(0.45), C_NAVY)
    add_textbox(slide, Inches(0.3), Inches(1.57), Inches(5.8), Inches(0.4),
                "Local (post-fix)", size_pt=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    headers_e = ["Metric", "HTTP Polling", "WebSocket", "Speedup"]
    data_local = [
        ["p50 (ms)",  "970",   ("180",  C_GREEN), ("5.4×", C_GREEN)],
        ["p95 (ms)",  "3 300", ("1 000",C_GREEN), ("3.3×", C_GREEN)],
        ["p99 (ms)",  "5 100", ("2 500",C_GREEN), ("2.0×", C_GREEN)],
        ["Average",   "1 236", ("328",  C_GREEN), ("3.8×", C_GREEN)],
    ]
    add_table(slide, Inches(0.3), Inches(2.05), Inches(5.8), 5, 4, headers_e, data_local,
              col_widths=[1.8, 1.6, 1.6, 1.2], font_size=13)

    # AWS
    add_rect(slide, Inches(6.7), Inches(1.55), Inches(6.1), Inches(0.45), C_TEAL)
    add_textbox(slide, Inches(6.7), Inches(1.57), Inches(6.1), Inches(0.4),
                "AWS ECS Fargate (post-fix)", size_pt=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    data_aws = [
        ["p50 (ms)",  "780",   ("~0",  C_GREEN),  ("—",      C_DGRAY)],
        ["p95 (ms)",  "2 500", ("62",  C_GREEN),  ("40×",  C_GREEN)],
        ["p99 (ms)",  "3 800", ("120", C_GREEN),  ("32×",  C_GREEN)],
        ["Average",   "942",   ("8",   C_GREEN),  ("116×", C_GREEN)],
    ]
    add_table(slide, Inches(6.7), Inches(2.05), Inches(6.1), 5, 4, headers_e, data_aws,
              col_widths=[1.8, 1.6, 1.6, 1.2], font_size=13)

    # Charts
    imgs = [
        (FIGS / "exp4"     / "exp4_latency_percentiles.png", Inches(0.3),  Inches(4.5), Inches(4.2), Inches(2.8), "Local – Percentiles"),
        (FIGS / "exp4"     / "exp4_avg_latency.png",          Inches(4.7),  Inches(4.5), Inches(4.2), Inches(2.8), "Local – Avg Latency"),
        (FIGS / "exp4_aws" / "exp4_latency_percentiles.png",  Inches(9.0),  Inches(4.5), Inches(4.1), Inches(2.8), "AWS – Percentiles"),
    ]
    for path, l, t, w2, h2, cap in imgs:
        add_image_safe(slide, path, l, t, w2, h2)
        add_textbox(slide, l, t + h2 + Inches(0.02), w2, Inches(0.3),
                    cap, size_pt=10, color=C_DGRAY, align=PP_ALIGN.CENTER)


def slide_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_rect(slide, 0, 0, Inches(0.45), H, C_TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.2), Inches(11), Inches(0.7),
                "Key Findings & Takeaways", size_pt=34, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.7), Inches(0.92), Inches(5.0), Inches(0.05), C_TEAL)

    findings = [
        ("Exp 1 — Scale-Out",
         "1→2 replicas give the biggest gain. Beyond 2, shared backends (Postgres pool, DynamoDB) become the bottleneck. PgBouncer would unlock further scaling."),
        ("Exp 2 — Partition Design",
         "LocalStack hides real storage effects; AWS validates theory: spreading partition keys reduces DynamoDB hot-partition throttling."),
        ("Exp 3 — Async Reactions",
         "SQS batch aggregation trades per-reaction latency for throughput efficiency and back-pressure resilience."),
        ("Exp 4 — WebSocket Fix",
         "Direct hub broadcast is essential for same-replica clients. Missing it sent every message through SNS→SQS (20 s delay!). Post-fix: WS is 116× faster than polling on AWS."),
    ]

    for idx, (title, body) in enumerate(findings):
        y = Inches(1.25) + idx * Inches(1.45)
        add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(1.3), C_NAVY)
        add_rect(slide, Inches(0.7), y, Inches(0.08), Inches(1.3), C_TEAL)
        add_textbox(slide, Inches(0.9), y + Inches(0.08), Inches(11.5), Inches(0.38),
                    title, size_pt=15, bold=True, color=C_TEAL2)
        add_textbox(slide, Inches(0.9), y + Inches(0.45), Inches(11.5), Inches(0.75),
                    body, size_pt=13, color=C_WHITE)


def slide_qa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, C_DARK)
    add_rect(slide, 0, 0, Inches(0.45), H, C_ORANGE)
    add_rect(slide, Inches(0.45), H - Inches(0.12), W - Inches(0.45), Inches(0.12), C_TEAL)

    add_textbox(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.4),
                "Thank You  /  Q & A", size_pt=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.5), Inches(3.5), Inches(8.0), Inches(0.06), C_TEAL)
    add_textbox(slide, Inches(1.2), Inches(3.7), Inches(11), Inches(0.6),
                "Figures:  live-chat-microservice/report/figures/",
                size_pt=16, color=C_TEAL2, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), Inches(4.4), Inches(11), Inches(0.5),
                "Load tests:  scripts/locustfile*.py  ·  run_experiment*.sh",
                size_pt=16, color=RGBColor(0xAA, 0xCC, 0xDD), align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Remove all default layouts (we always use layout[6] = blank)
    slide_title(prs)
    slide_agenda(prs)
    slide_architecture(prs)
    slide_tech_stack(prs)
    slide_exp1_overview(prs)
    slide_exp1_charts(prs)
    slide_exp2_overview(prs)
    slide_exp2_charts(prs)
    slide_exp3(prs)
    slide_exp4_bug(prs)
    slide_exp4_results(prs)
    slide_findings(prs)
    slide_qa(prs)

    out = Path(__file__).resolve().parent.parent / "docs" / "LiveChat_Presentation.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"✅  Saved → {out}")


if __name__ == "__main__":
    main()
